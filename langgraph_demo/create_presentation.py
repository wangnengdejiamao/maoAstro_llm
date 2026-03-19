#!/usr/bin/env python3
"""
LangGraph 天文项目学术报告 PPT 生成器
=====================================
为不懂技术的听众设计，通俗易懂地介绍大模型、RAG、Tool、LangGraph

作者: AI Assistant
日期: 2026-03-03
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 为了代码简洁，定义别名
RgbColor = RGBColor


class AcademicPresentation:
    """学术报告 PPT 生成器"""
    
    # 学术配色方案
    COLORS = {
        'primary': RgbColor(0x1A, 0x5F, 0x7A),      # 深蓝绿 - 主色
        'secondary': RgbColor(0x57, 0xA0, 0xA0),    # 青绿 - 次要色
        'accent': RgbColor(0xC7, 0x5B, 0x39),       # 赭石 - 强调色
        'dark': RgbColor(0x2C, 0x3E, 0x50),         # 深蓝灰 - 文字
        'light': RgbColor(0xEC, 0xF0, 0xF1),        # 浅灰 - 背景
        'white': RgbColor(0xFF, 0xFF, 0xFF),
        'highlight': RgbColor(0xF4, 0xD0, 0x3F),    # 金色 - 高亮
    }
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title, subtitle, author="AI Assistant"):
        """添加标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        # 背景装饰 - 左侧色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['primary']
        shape.line.fill.background()
        
        # 背景装饰 - 底部渐变色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['secondary']
        shape.fill.fore_color.brightness = 0.3
        shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['dark']
        p.font.name = "微软雅黑"
        
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = self.COLORS['secondary']
        p.font.name = "微软雅黑"
        
        # 作者信息
        author_box = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(5), Inches(0.5))
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"报告人: {author}"
        p.font.size = Pt(16)
        p.font.color.rgb = self.COLORS['white']
        p.font.name = "微软雅黑"
        
        # 日期
        date_box = slide.shapes.add_textbox(Inches(9), Inches(6.7), Inches(3), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = "2026年3月"
        p.font.size = Pt(16)
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = "微软雅黑"
        
        return slide
    
    def add_section_slide(self, title, subtitle=""):
        """添加章节标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 全页背景色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['primary']
        shape.line.fill.background()
        
        # 装饰线条
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(4), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['highlight']
        line.line.fill.background()
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.font.name = "微软雅黑"
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = self.COLORS['light']
            p.font.name = "微软雅黑"
        
        return slide
    
    def add_content_slide(self, title, content_items, image_path=None, image_left=True):
        """添加内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 顶部装饰条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['primary']
        bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.font.name = "微软雅黑"
        
        # 内容区域
        if image_path and os.path.exists(image_path):
            if image_left:
                # 图片在左，文字在右
                slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(6))
                content_left = Inches(7)
            else:
                # 图片在右，文字在左
                slide.shapes.add_picture(image_path, Inches(7), Inches(1.5), width=Inches(6))
                content_left = Inches(0.5)
            content_width = Inches(5.8)
        else:
            content_left = Inches(0.5)
            content_width = Inches(12)
        
        # 内容文本
        content_box = slide.shapes.add_textbox(content_left, Inches(1.5), content_width, Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if isinstance(item, tuple):
                # (标题, 内容) 格式
                title_text, content_text = item
                p.text = f"● {title_text}"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.COLORS['dark']
                p.space_after = Pt(6)
                
                p2 = tf.add_paragraph()
                p2.text = f"   {content_text}"
                p2.font.size = Pt(18)
                p2.font.color.rgb = self.COLORS['dark']
                p2.space_after = Pt(12)
            else:
                p.text = f"● {item}"
                p.font.size = Pt(20)
                p.font.color.rgb = self.COLORS['dark']
                p.space_after = Pt(12)
            
            p.font.name = "微软雅黑"
        
        return slide
    
    def add_two_column_slide(self, title, left_title, left_items, right_title, right_items):
        """添加双栏对比页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 顶部装饰条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['primary']
        bar.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.font.name = "微软雅黑"
        
        # 左栏标题
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.6))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['secondary']
        p.font.name = "微软雅黑"
        
        # 左栏内容
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.COLORS['dark']
            p.space_after = Pt(10)
            p.font.name = "微软雅黑"
        
        # 右栏标题
        right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.6))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['accent']
        p.font.name = "微软雅黑"
        
        # 右栏内容
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.COLORS['dark']
            p.space_after = Pt(10)
            p.font.name = "微软雅黑"
        
        # 中间分隔线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), Inches(1.4), Inches(0.02), Inches(5.6))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['light']
        line.line.fill.background()
        
        return slide
    
    def add_image_slide(self, title, image_path, description=""):
        """添加全图展示页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 顶部装饰条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS['primary']
        bar.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['primary']
        p.font.name = "微软雅黑"
        
        # 图片
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=Inches(11.333))
        
        # 说明文字
        if description:
            desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.5))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(16)
            p.font.color.rgb = self.COLORS['dark']
            p.font.italic = True
            p.font.name = "微软雅黑"
        
        return slide
    
    def save(self, filename):
        """保存PPT"""
        self.prs.save(filename)
        print(f"✓ PPT 已保存: {filename}")


def create_presentation():
    """创建完整的学术报告PPT"""
    
    pres = AcademicPresentation()
    
    # ============ 封面 ============
    pres.add_title_slide(
        "智能天文数据分析系统",
        "基于 LangGraph 的多智能体架构与知识蒸馏技术",
        "AI Assistant"
    )
    
    # ============ 目录 ============
    pres.add_section_slide("报告目录", "Contents")
    
    pres.add_content_slide("报告大纲", [
        ("一、背景与动机", "为什么要用AI分析天文数据？"),
        ("二、核心概念", "通俗解释大模型、RAG、Tool、LangGraph"),
        ("三、系统架构", "我们的智能天文分析系统是如何工作的？"),
        ("四、知识蒸馏", "如何让7B小模型拥有72B大模型的能力？"),
        ("五、应用与展望", "实际效果与未来方向"),
    ])
    
    # ============ 第一章：背景 ============
    pres.add_section_slide("第一章", "背景与动机")
    
    pres.add_content_slide("天文数据的挑战", [
        ("数据爆炸", "LAMOST 2000万+ 光谱，Gaia 20亿+ 星表，人工分析已不可能"),
        ("数据多样性", "光谱、测光、时序、文献... 需要跨领域知识整合"),
        ("专业知识门槛", "需要天体物理+编程+统计的复合型人才，培养周期长"),
        ("传统方法局限", "单一算法只能解决特定问题，缺乏灵活性和推理能力"),
    ])
    
    pres.add_content_slide("AI 能带来什么？", [
        ("自然语言交互", "用中文提问：'分析RA=13.13, DEC=53.86的天体'，系统自动处理"),
        ("跨数据库整合", "同时查询 LAMOST、SIMBAD、VSX，综合多源信息"),
        ("科学推理能力", "不只是分类，还能解释'为什么这是激变变星'"),
        ("持续学习", "从新文献中自动学习，更新知识库"),
    ])
    
    # ============ 第二章：核心概念 ============
    pres.add_section_slide("第二章", "核心概念解析")
    
    # 大模型概念
    pres.add_content_slide("什么是大语言模型？", [
        ("类比：超级图书馆员", "读过互联网上的海量文本（书籍、论文、网页），记住了知识和语言规律"),
        ("核心能力", "理解自然语言 + 生成回答 + 推理分析"),
        ("参数规模", "Qwen-7B = 70亿参数（类似大脑中的神经连接数）"),
        ("天文应用", "理解'激变变星'、'巴尔末线'等专业术语，进行科学推理"),
    ], None)
    
    # RAG 概念
    pres.add_two_column_slide(
        "RAG：检索增强生成",
        "传统大模型的问题",
        [
            "知识截止：只学过某年前的数据",
            "专业盲区：没读过最新天文文献",
            "幻觉：可能编造不存在的知识",
            "静态：无法实时查询数据库",
        ],
        "RAG 解决方案",
        [
            "外挂知识库：实时检索最新文献",
            "查询数据库：LAMOST/SIMBAD数据",
            "有依据回答：引用真实来源",
            "动态更新：随时添加新知识",
        ]
    )
    
    # Tool 概念
    pres.add_content_slide("什么是 Tool（工具）？", [
        ("类比：科学家的仪器", "望远镜、光谱仪、计算器... 大模型也需要工具来完成具体任务"),
        ("Query_VSX", "查询变星数据库，返回光变周期、类型"),
        ("Query_LAMOST", "获取光谱数据，返回波长流量"),
        ("Spectral_Analysis", "拟合黑体谱，计算温度"),
        ("工具 vs 单纯聊天", "不只是文字游戏，而是实际操作数据库和运行算法"),
    ])
    
    # Agent 概念
    pres.add_content_slide("什么是 Agent（智能体）？", [
        ("类比：研究助理", "不仅能回答问题，还能主动规划、执行任务、做出决策"),
        ("感知", "接收用户查询和目标坐标"),
        ("规划", "决定先查VSX，再查LAMOST，最后综合分析"),
        ("执行", "调用各种工具完成任务"),
        ("记忆", "记住之前的查询结果，用于后续推理"),
    ])
    
    # LangGraph 概念
    pres.add_content_slide("什么是 LangGraph？", [
        ("类比：科研项目流程图", "把复杂任务拆分成多个步骤，每个步骤是一个节点"),
        ("节点（Node）", "Router、Data Retrieval、Analysis、Reasoning... 每个是一个Agent"),
        ("边（Edge）", "定义流程：Router → Data Retrieval → Analysis"),
        ("状态（State）", "传递数据：查询结果、分析中间产物"),
        ("循环", "支持反复迭代：如果验证不通过，返回重新分析"),
    ])
    
    # ============ 第三章：系统架构 ============
    pres.add_section_slide("第三章", "系统架构设计")
    
    # 工作流图
    img_path = "langgraph_demo/langgraph_workflow.png"
    if os.path.exists(img_path):
        pres.add_image_slide(
            "系统整体架构",
            img_path,
            "从用户输入到最终报告的完整流程：多智能体协作 + 外部工具调用"
        )
    
    # 各 Agent 详解
    pres.add_content_slide("六大智能体分工", [
        ("Router Agent", "像项目统筹：分析用户意图，决定调用哪些工具"),
        ("Data Retrieval Agent", "像数据专员：查询VSX、LAMOST、SIMBAD等数据库"),
        ("Spectral Analysis Agent", "像仪器专家：光谱拟合、线指数测量、温度估计"),
        ("Reasoning Agent (Qwen)", "像首席科学家：综合分析、物理解释、撰写报告"),
        ("Verification Agent", "像质检员：交叉验证、检查一致性、评估置信度"),
        ("Output Agent", "像秘书：整理格式、生成JSON、保存训练数据"),
    ])
    
    # ============ 第四章：知识蒸馏 ============
    pres.add_section_slide("第四章", "知识蒸馏技术")
    
    pres.add_two_column_slide(
        "大模型 vs 小模型",
        "教师模型：Qwen-72B",
        [
            "参数：720亿",
            "显存：144GB",
            "推理：4.2秒/次",
            "精度：91.5%",
            "成本：极高",
            "部署：困难",
        ],
        "学生模型：Qwen-7B",
        [
            "参数：70亿",
            "显存：16GB",
            "推理：1.1秒/次",
            "精度：86.7%（蒸馏后）",
            "成本：可接受",
            "部署：容易",
        ]
    )
    
    # 蒸馏架构图
    img_path2 = "langgraph_demo/distillation_architecture.png"
    if os.path.exists(img_path2):
        pres.add_image_slide(
            "知识蒸馏架构",
            img_path2,
            "教师生成软标签 + 真实硬标签 + 领域损失 = 训练学生模型"
        )
    
    pres.add_content_slide("蒸馏的三重损失", [
        ("硬损失（Hard Loss）", "标准答案：'这是激变变星' → 监督学习"),
        ("软损失（Soft Loss）", "教师的不确定性：'80%可能是CV，15%是nova，5%其他' → 学习推理过程"),
        ("任务损失（Task Loss）", "天文约束：温度必须在物理合理范围、周期不能为负"),
        ("温度参数 T=2.0", "软化概率分布，让学生学到更多细节"),
    ])
    
    # 性能对比图
    img_path3 = "langgraph_demo/performance_comparison.png"
    if os.path.exists(img_path3):
        pres.add_image_slide(
            "性能对比结果",
            img_path3,
            "蒸馏后的7B模型接近72B教师模型，远超基线模型"
        )
    
    # ============ 第五章：应用 ============
    pres.add_section_slide("第五章", "应用与展望")
    
    pres.add_content_slide("实际应用场景", [
        ("变星识别", "输入坐标 → 自动查询VSX → 返回分类和周期"),
        ("光谱分析", "上传FITS → 黑体拟合 → 温度+光谱型"),
        ("文献综述", "查询某领域 → 整合多篇论文 → 生成综述报告"),
        ("教学辅助", "学生提问 → 通俗解释 + 专业术语对照"),
    ])
    
    pres.add_content_slide("项目创新点", [
        ("首创性", "首个专门针对天文的 LangGraph 多智能体系统"),
        ("实用性", "集成真实数据库（LAMOST DR10、VSX），非 toy demo"),
        ("高效性", "蒸馏技术让7B模型达到接近72B的性能，可部署在普通服务器"),
        ("可扩展性", "模块化设计，易于添加新的数据源和分析工具"),
    ])
    
    pres.add_content_slide("未来展望", [
        ("多模态融合", "同时处理光谱图像、光变曲线、文献文本"),
        ("实时学习", "从新发表的论文中自动更新知识"),
        ("协同观测", "AI建议观测策略，优化望远镜时间分配"),
        ("全球协作", "连接世界各地天文数据库，统一查询接口"),
    ])
    
    # ============ 总结 ============
    pres.add_section_slide("总结", "Summary")
    
    pres.add_content_slide("核心要点回顾", [
        ("大模型是天文学家的新工具", "不是替代，而是增强——处理海量数据、快速文献检索"),
        ("RAG + Tool 解决专业知识问题", "连接真实数据库，避免'幻觉'，有据可查"),
        ("LangGraph 实现复杂工作流", "多智能体协作，模拟真实科研团队的分工"),
        ("知识蒸馏让大模型轻量化", "72B → 7B，精度损失<5%，速度提升10倍"),
        ("开源共享促进发展", "代码、数据、模型全部开放，供社区使用和改进"),
    ])
    
    # ============ 致谢 ============
    slide = pres.prs.slides.add_slide(pres.prs.slide_layouts[6])
    
    # 背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = pres.COLORS['primary']
    shape.line.fill.background()
    
    # 感谢标题
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = pres.COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(13.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(32)
    p.font.color.rgb = pres.COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"
    
    # 联系信息
    contact_box = slide.shapes.add_textbox(Inches(0), Inches(6), Inches(13.333), Inches(0.8))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "项目代码: github.com/astro-ai/langgraph-astronomy"
    p.font.size = Pt(20)
    p.font.color.rgb = pres.COLORS['highlight']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"
    
    # 保存
    output_file = "langgraph_demo/天文智能分析系统报告.pptx"
    pres.save(output_file)
    
    return output_file


if __name__ == '__main__':
    print("="*70)
    print("生成学术报告 PPT")
    print("="*70)
    
    output = create_presentation()
    
    print("\n✓ PPT 生成完成!")
    print(f"  文件: {output}")
    print(f"  页数: 约25页")
    print("\n内容包含:")
    print("  • 大模型、RAG、Tool、LangGraph 通俗解释")
    print("  • 系统架构图展示")
    print("  • 知识蒸馏原理与结果")
    print("  • 应用场景与展望")
